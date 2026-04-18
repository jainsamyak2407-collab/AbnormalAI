"""
Unit tests for the PPTX renderer and chart renderer.
These tests use minimal fixtures — no AI calls, no metrics pipeline.
"""

from __future__ import annotations

import io

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from export.chart_renderer import render_chart
from export.pptx_renderer import render_presentation
from export.slide_theme import SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN


# ── fixtures ──────────────────────────────────────────────────────────────────

FOOTER = "Meridian Healthcare · Q1 2026 · CISO View · Slide {n} of 5"


def _title_slide(n: int = 1) -> dict:
    return {
        "slide_number": n,
        "slide_type": "title",
        "headline": "Meridian Healthcare",
        "subtitle": "EXECUTIVE BRIEF · Q1 2026 · Prepared for the Board of Directors",
        "footer": FOOTER.format(n=n),
        "evidence_refs": [],
    }


def _thesis_slide(n: int = 2) -> dict:
    return {
        "slide_number": n,
        "slide_type": "thesis",
        "thesis_sentence": "Meridian cut VIP inbox exposure by 67% while closing its user reporting gap, with auto-remediation and acquisition hygiene requiring Q2 investment.",
        "thesis_tagline": "Protection Holds. Gaps Linger.",
        "footer": FOOTER.format(n=n),
        "evidence_refs": [],
    }


def _what_happened_slide(n: int = 3) -> dict:
    return {
        "slide_number": n,
        "slide_type": "what_happened",
        "headline": "User reporting crossed 40% in March while VIP attacks fell to one per month",
        "chart": {
            "type": "trend_line",
            "title": "VIP Inbox Attacks by Month",
            "caption": "Monthly count of threats reaching VIP inboxes",
            "source_note": "threats.csv",
            "data": {
                "x_labels": ["Jan", "Feb", "Mar"],
                "series": [{"label": "VIP attacks", "values": [3, 1, 1]}],
                "target": None,
            },
            "evidence_refs": ["E1"],
        },
        "callouts": [
            {"number": "1,847", "label": "Threats blocked in Q1", "context": "Up 12% quarter-over-quarter.", "color": "ink"},
            {"number": "45%", "label": "User reporting rate in March", "context": "Crossed the 40% board success criterion for the first time.", "color": "success"},
            {"number": "2.3 min", "label": "Median MTTR", "context": "Beats the industry p75 of 7.4 minutes.", "color": "success"},
        ],
        "footer": FOOTER.format(n=n),
        "evidence_refs": ["E1", "E3", "E9"],
    }


def _what_needs_attention_slide(n: int = 4) -> dict:
    return {
        "slide_number": n,
        "slide_type": "what_needs_attention",
        "headline": "Auto-remediation gap and credential submissions are the two compounding exposures",
        "chart": {
            "type": "benchmark_bars",
            "title": "Meridian vs. Healthcare Peers",
            "caption": "Key KPIs vs industry p25/p50/p75",
            "source_note": "benchmarks",
            "data": {
                "metrics": [
                    {"label": "Auto-remediation rate", "value": 67.6, "p25": 60.0, "p50": 75.0, "p75": 85.0, "unit": "%", "higher_is_better": True},
                    {"label": "Median MTTR", "value": 2.3, "p25": 12.0, "p50": 7.4, "p75": 4.2, "unit": "min", "higher_is_better": False},
                ]
            },
            "evidence_refs": ["E12"],
        },
        "callouts": [
            {"number": "67.6%", "label": "Auto-remediation rate", "context": "Below the healthcare peer median of 75.0%, driving 20 manual SOC interventions.", "color": "accent"},
            {"number": "6.4%", "label": "Credential submission rate in March", "context": "Tripled from 2.1% in January, above the 3.0% success threshold.", "color": "accent"},
            {"number": "72.1%", "label": "T-002 posture pass rate", "context": "7.7 points below T-001, the clearest acquisition hygiene gap.", "color": "warning"},
        ],
        "footer": FOOTER.format(n=n),
        "evidence_refs": ["E12", "E7", "E14"],
    }


def _ask_slide(n: int = 5) -> dict:
    return {
        "slide_number": n,
        "slide_type": "the_ask",
        "headline": "Three actions close the compounding risk before Q2 board review",
        "recommendations": [
            {
                "kind": "BUDGET",
                "headline": "Fund the auto-remediation tuning sprint before April 30 to close the 7.4-point gap",
                "rationale": "Auto-remediation at 67.6% generates 20 manual SOC interventions per quarter; reaching 75% eliminates them.",
                "evidence_refs": ["E12"],
            },
            {
                "kind": "POLICY",
                "headline": "Mandate MFA enforcement on T-001 via policy update this month",
                "rationale": "13 consecutive weekly failures with up to 70 affected users represent the highest unresolved risk in posture.",
                "evidence_refs": ["E14"],
            },
            {
                "kind": "TRAINING",
                "headline": "Deploy credential-submission training for the 24 highest-risk users by end of April",
                "rationale": "Submission rate tripled from 2.1% to 6.4% in Q1; targeted training brings it below 3.0% before Q2 close.",
                "evidence_refs": ["E7"],
            },
        ],
        "closing_ask": "To close the auto-remediation gap and complete the T-002 integration, we are requesting budget approval and a policy mandate at the April board meeting.",
        "footer": FOOTER.format(n=n),
        "evidence_refs": ["E12", "E14", "E7"],
    }


def _minimal_deck() -> list[dict]:
    return [
        _title_slide(),
        _thesis_slide(),
        _what_happened_slide(),
        _what_needs_attention_slide(),
        _ask_slide(),
    ]


# ── chart renderer tests ──────────────────────────────────────────────────────

class TestChartRenderer:
    def test_trend_line_returns_png(self):
        spec = {
            "type": "trend_line",
            "data": {
                "x_labels": ["Jan", "Feb", "Mar"],
                "series": [{"label": "VIP attacks", "values": [3, 1, 1]}],
                "target": None,
            },
        }
        png = render_chart(spec, dpi=72)
        assert isinstance(png, bytes)
        assert png[:4] == b"\x89PNG"  # PNG magic bytes

    def test_benchmark_bars_returns_png(self):
        spec = {
            "type": "benchmark_bars",
            "data": {
                "metrics": [
                    {"label": "Auto-remediation", "value": 67.6, "p25": 60.0, "p50": 75.0, "p75": 85.0, "unit": "%", "higher_is_better": True},
                ]
            },
        }
        png = render_chart(spec, dpi=72)
        assert png[:4] == b"\x89PNG"

    def test_department_bars_returns_png(self):
        spec = {
            "type": "department_bars",
            "data": {
                "categories": [{"label": "Legal", "value": 19.7}, {"label": "Engineering", "value": 42.1}],
                "threshold": 40.0,
                "threshold_label": "Success criterion",
            },
        }
        png = render_chart(spec, dpi=72)
        assert png[:4] == b"\x89PNG"

    def test_criteria_scorecard_returns_png(self):
        spec = {
            "type": "criteria_scorecard",
            "data": {
                "rows": [
                    {"criterion": "VIP attacks < 1/month", "target": "1", "actual": "1", "met": True},
                    {"criterion": "Reporting rate > 40%", "target": "40%", "actual": "45%", "met": True},
                ]
            },
        }
        png = render_chart(spec, dpi=72)
        assert png[:4] == b"\x89PNG"

    def test_empty_data_returns_placeholder(self):
        spec = {"type": "trend_line", "data": {}}
        png = render_chart(spec, dpi=72)
        assert png[:4] == b"\x89PNG"

    def test_unknown_type_returns_placeholder(self):
        spec = {"type": "unknown_chart", "data": {}}
        png = render_chart(spec, dpi=72)
        assert png[:4] == b"\x89PNG"


# ── pptx renderer tests ───────────────────────────────────────────────────────

class TestPptxRenderer:
    def _render(self, slides: list[dict] | None = None, charts: dict | None = None) -> PptxPresentation:
        if slides is None:
            slides = _minimal_deck()
        if charts is None:
            # Render real chart PNGs for slides 3 and 4
            charts = {}
            for s in slides:
                chart_spec = s.get("chart")
                if chart_spec:
                    from export.chart_renderer import render_chart
                    charts[s["slide_number"]] = render_chart(chart_spec, dpi=72)
        raw = render_presentation(slides, charts)
        assert isinstance(raw, bytes), "render_presentation must return bytes"
        return PptxPresentation(io.BytesIO(raw))

    def test_output_is_pptx_bytes(self):
        slides = _minimal_deck()
        raw = render_presentation(slides, {})
        assert isinstance(raw, bytes)
        assert len(raw) > 1000

    def test_exactly_5_slides(self):
        prs = self._render()
        assert len(prs.slides) == 5

    def test_correct_dimensions(self):
        prs = self._render()
        assert abs(prs.slide_width.inches - SLIDE_WIDTH_IN) < 0.01
        assert abs(prs.slide_height.inches - SLIDE_HEIGHT_IN) < 0.01

    def test_all_slide_types_render(self):
        # Each slide type must not raise
        for slide_type in ["title", "thesis", "what_happened", "what_needs_attention", "the_ask"]:
            slides = [s for s in _minimal_deck() if s["slide_type"] == slide_type]
            raw = render_presentation(slides, {})
            assert len(raw) > 500

    def test_missing_chart_png_renders_gracefully(self):
        # Slide 3 and 4 have charts but we pass no PNGs
        prs = self._render(charts={})
        assert len(prs.slides) == 5

    def test_unknown_slide_type_renders_gracefully(self):
        slides = _minimal_deck()
        slides[0]["slide_type"] = "not_a_real_type"
        raw = render_presentation(slides, {})
        prs = PptxPresentation(io.BytesIO(raw))
        assert len(prs.slides) == 5
