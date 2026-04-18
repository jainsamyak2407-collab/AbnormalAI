"""
pptx_renderer.py — compiles 5 SlideContent dicts into a .pptx file via python-pptx.
All layout, color, and typography comes from slide_theme.py.
AI never touches this layer.
"""

from __future__ import annotations

import io
import logging
from typing import Any

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

from export.slide_theme import (
    SLIDE_THEME, FONTS,
    SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN,
    MARGIN_TOP_IN, MARGIN_LR_IN, MARGIN_BOT_IN,
    ACCENT_RULE_X_IN, ACCENT_RULE_Y_IN, ACCENT_RULE_W_IN, ACCENT_RULE_H_IN,
    FOOTER_Y_IN, FOOTER_H_IN, FOOTER_FONT_PT,
    DARK_BAND_H_IN,
)

logger = logging.getLogger(__name__)


# ── colour helpers ────────────────────────────────────────────────────────────

def _rgb(key: str) -> RGBColor:
    r, g, b = SLIDE_THEME[key]
    return RGBColor(r, g, b)


# ── shape helpers ─────────────────────────────────────────────────────────────

def _add_rect(slide, x: float, y: float, w: float, h: float,
              fill_key: str, line: bool = False) -> Any:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_key)
    if not line:
        shape.line.fill.background()
    return shape


def _add_text(
    slide,
    text: str,
    x: float, y: float, w: float, h: float,
    font_name: str,
    font_pt: float,
    color_key: str,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
) -> Any:
    txbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.name = font_name
    run.font.size = Pt(font_pt)
    run.font.color.rgb = _rgb(color_key)
    run.font.bold = bold
    return txbox


def _add_footer(slide, footer_text: str, dark_bg: bool = False) -> None:
    color_key = "ink_on_dark" if dark_bg else "ink_secondary"
    _add_text(
        slide, footer_text,
        x=MARGIN_LR_IN,
        y=FOOTER_Y_IN,
        w=SLIDE_WIDTH_IN - MARGIN_LR_IN * 2,
        h=FOOTER_H_IN,
        font_name=FONTS["mono"],
        font_pt=FOOTER_FONT_PT,
        color_key=color_key,
        align=PP_ALIGN.RIGHT,
    )


def _add_accent_rule(slide) -> None:
    _add_rect(slide,
              x=ACCENT_RULE_X_IN, y=ACCENT_RULE_Y_IN,
              w=ACCENT_RULE_W_IN, h=ACCENT_RULE_H_IN,
              fill_key="accent")


def _add_chart_image(slide, png_bytes: bytes, x: float, y: float, w: float, h: float) -> None:
    img_stream = io.BytesIO(png_bytes)
    slide.shapes.add_picture(img_stream, Inches(x), Inches(y), Inches(w), Inches(h))


# ── full-bleed background ─────────────────────────────────────────────────────

def _fill_background(slide, fill_key: str) -> None:
    _add_rect(slide, 0, 0, SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN, fill_key)


# ── slide builders ────────────────────────────────────────────────────────────

def _build_title_slide(slide, content: dict) -> None:
    _fill_background(slide, "bg_alt")

    # Coral accent rule
    _add_accent_rule(slide)

    # Label (EXECUTIVE BRIEF / QUARTERLY BUSINESS REVIEW)
    subtitle: str = content.get("subtitle") or ""
    label = subtitle.split("·")[0].strip() if "·" in subtitle else subtitle
    _add_text(slide, label,
              x=MARGIN_LR_IN, y=1.2, w=8.0, h=0.4,
              font_name=FONTS["mono"], font_pt=11,
              color_key="accent", bold=False, align=PP_ALIGN.LEFT)

    # Customer name — hero
    customer: str = content.get("headline") or ""
    _add_text(slide, customer,
              x=MARGIN_LR_IN, y=1.75, w=9.5, h=1.6,
              font_name=FONTS["serif"], font_pt=52,
              color_key="ink_on_dark", bold=False, align=PP_ALIGN.LEFT)

    # Period / prepared-for
    period_line = subtitle.split("·", 1)[1].strip() if "·" in subtitle else ""
    if period_line:
        _add_text(slide, period_line,
                  x=MARGIN_LR_IN, y=3.5, w=8.0, h=0.4,
                  font_name=FONTS["sans"], font_pt=15,
                  color_key="ink_on_dark", align=PP_ALIGN.LEFT)

    # Attribution footer
    _add_footer(slide, content.get("footer", ""), dark_bg=True)


def _build_thesis_slide(slide, content: dict) -> None:
    _fill_background(slide, "bg_alt")
    _add_accent_rule(slide)

    thesis: str = content.get("thesis_sentence") or ""
    _add_text(slide, thesis,
              x=MARGIN_LR_IN, y=1.6, w=11.5, h=2.8,
              font_name=FONTS["serif"], font_pt=40,
              color_key="ink_on_dark", align=PP_ALIGN.LEFT)

    tagline: str = content.get("thesis_tagline") or ""
    if tagline:
        _add_text(slide, tagline,
                  x=MARGIN_LR_IN, y=4.6, w=9.0, h=0.6,
                  font_name=FONTS["sans"], font_pt=14,
                  color_key="ink_secondary", align=PP_ALIGN.LEFT)

    _add_footer(slide, content.get("footer", ""), dark_bg=True)


def _build_callout_stack(slide, callouts: list[dict],
                         x: float, y_start: float, w: float) -> None:
    """Render 3 callouts stacked vertically."""
    slot_h = 1.55
    for i, c in enumerate(callouts[:3]):
        y = y_start + i * slot_h

        # Hairline separator (except before first)
        if i > 0:
            _add_rect(slide, x, y - 0.05, w, 0.008, fill_key="rule")

        number_str: str = c.get("number", "")
        label_str: str = c.get("label", "")
        context_str: str = c.get("context", "")
        color_key = {
            "accent": "accent",
            "success": "success",
            "warning": "warning",
        }.get(c.get("color", "ink"), "ink_primary")

        _add_text(slide, number_str,
                  x=x, y=y, w=w, h=0.6,
                  font_name=FONTS["serif"], font_pt=30,
                  color_key=color_key, bold=False)

        _add_text(slide, label_str,
                  x=x, y=y + 0.58, w=w, h=0.35,
                  font_name=FONTS["sans"], font_pt=12,
                  color_key="ink_primary")

        _add_text(slide, context_str,
                  x=x, y=y + 0.92, w=w, h=0.55,
                  font_name=FONTS["sans"], font_pt=10,
                  color_key="ink_secondary")


def _build_data_slide(slide, content: dict, chart_png: bytes | None,
                      chart_left: bool) -> None:
    """Shared builder for slides 3 (chart left) and 4 (chart right)."""
    _fill_background(slide, "bg_primary")
    _add_accent_rule(slide)

    headline: str = content.get("headline") or ""
    _add_text(slide, headline,
              x=MARGIN_LR_IN, y=ACCENT_RULE_Y_IN + 0.15,
              w=SLIDE_WIDTH_IN - MARGIN_LR_IN * 2, h=0.65,
              font_name=FONTS["serif"], font_pt=26,
              color_key="ink_primary", bold=False)

    chart_w = 7.8
    callout_w = 4.3
    content_y = 1.3
    content_h = 5.4

    if chart_left:
        chart_x = MARGIN_LR_IN
        callout_x = MARGIN_LR_IN + chart_w + 0.25
    else:
        callout_x = MARGIN_LR_IN
        chart_x = MARGIN_LR_IN + callout_w + 0.25

    if chart_png:
        _add_chart_image(slide, chart_png,
                         x=chart_x, y=content_y,
                         w=chart_w, h=content_h - 0.1)
    else:
        # Placeholder box if chart unavailable
        _add_rect(slide, chart_x, content_y, chart_w, content_h - 0.1, fill_key="rule")

    callouts = content.get("callouts") or []
    _build_callout_stack(slide, callouts, callout_x, content_y + 0.2, callout_w - 0.1)

    _add_footer(slide, content.get("footer", ""))


def _build_ask_slide(slide, content: dict) -> None:
    _fill_background(slide, "bg_primary")
    _add_accent_rule(slide)

    headline: str = content.get("headline") or ""
    _add_text(slide, headline,
              x=MARGIN_LR_IN, y=ACCENT_RULE_Y_IN + 0.15,
              w=SLIDE_WIDTH_IN - MARGIN_LR_IN * 2, h=0.65,
              font_name=FONTS["serif"], font_pt=24,
              color_key="ink_primary", bold=False)

    recs = content.get("recommendations") or []
    rec_slot_h = 1.35
    rec_y_start = 1.3

    KIND_COLORS = {
        "POLICY": "ink_secondary",
        "BUDGET": "accent",
        "HEADCOUNT": "warning",
        "EXPANSION": "success",
        "TRAINING": "ink_secondary",
        "RENEWAL": "success",
    }

    for i, rec in enumerate(recs[:3]):
        y = rec_y_start + i * rec_slot_h
        kind: str = rec.get("kind", "")
        rec_headline: str = rec.get("headline", "")
        rationale: str = rec.get("rationale", "")

        # Kind chip (small rect + label)
        chip_w = 1.1
        chip_x = SLIDE_WIDTH_IN - MARGIN_LR_IN - chip_w
        _add_rect(slide, chip_x, y + 0.08, chip_w, 0.28, fill_key="rule")
        _add_text(slide, kind,
                  x=chip_x + 0.05, y=y + 0.08, w=chip_x, h=0.28,
                  font_name=FONTS["mono"], font_pt=8,
                  color_key=KIND_COLORS.get(kind, "ink_secondary"),
                  align=PP_ALIGN.LEFT)

        # Number label
        _add_text(slide, f"{i + 1}.",
                  x=MARGIN_LR_IN, y=y, w=0.35, h=0.5,
                  font_name=FONTS["serif"], font_pt=18,
                  color_key="accent")

        _add_text(slide, rec_headline,
                  x=MARGIN_LR_IN + 0.4, y=y, w=SLIDE_WIDTH_IN - MARGIN_LR_IN * 2 - chip_w - 0.6, h=0.5,
                  font_name=FONTS["serif"], font_pt=16,
                  color_key="ink_primary")

        _add_text(slide, rationale,
                  x=MARGIN_LR_IN + 0.4, y=y + 0.48, w=SLIDE_WIDTH_IN - MARGIN_LR_IN * 2 - 0.6, h=0.7,
                  font_name=FONTS["sans"], font_pt=11,
                  color_key="ink_secondary")

        # Row separator
        if i < 2:
            _add_rect(slide, MARGIN_LR_IN, y + rec_slot_h - 0.05,
                      SLIDE_WIDTH_IN - MARGIN_LR_IN * 2, 0.006, fill_key="rule")

    # Dark footer band
    band_y = SLIDE_HEIGHT_IN - DARK_BAND_H_IN
    _add_rect(slide, 0, band_y, SLIDE_WIDTH_IN, DARK_BAND_H_IN, fill_key="bg_alt")

    closing_ask: str = content.get("closing_ask") or ""
    if closing_ask:
        _add_text(slide, closing_ask,
                  x=MARGIN_LR_IN, y=band_y + 0.35,
                  w=SLIDE_WIDTH_IN - MARGIN_LR_IN * 2, h=DARK_BAND_H_IN - 0.5,
                  font_name=FONTS["serif"], font_pt=16,
                  color_key="ink_on_dark", align=PP_ALIGN.LEFT)

    _add_footer(slide, content.get("footer", ""))


# ── public API ────────────────────────────────────────────────────────────────

def render_presentation(slides: list[dict], chart_pngs: dict[int, bytes]) -> bytes:
    """
    Compile 5 SlideContent dicts into a .pptx and return bytes.

    slides: list of SlideContent-compatible dicts (slide_number, slide_type, ...)
    chart_pngs: {slide_number: png_bytes} — only slides 3 and 4 use a chart
    """
    prs = PptxPresentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    blank_layout = prs.slide_layouts[6]  # "Blank" layout

    # Sort slides defensively
    ordered = sorted(slides, key=lambda s: s.get("slide_number", 0))

    for content in ordered:
        slide = prs.slides.add_slide(blank_layout)
        sn: int = content.get("slide_number", 0)
        st: str = content.get("slide_type", "title")
        png = chart_pngs.get(sn)

        try:
            if st == "title":
                _build_title_slide(slide, content)
            elif st == "thesis":
                _build_thesis_slide(slide, content)
            elif st == "what_happened":
                _build_data_slide(slide, content, png, chart_left=True)
            elif st == "what_needs_attention":
                _build_data_slide(slide, content, png, chart_left=False)
            elif st == "the_ask":
                _build_ask_slide(slide, content)
            else:
                logger.warning("Unknown slide type '%s' for slide %d; leaving blank.", st, sn)
        except Exception as exc:
            logger.error("Slide %d (%s) render failed: %s", sn, st, exc)
            # Leave the slide blank rather than crashing the whole deck

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
